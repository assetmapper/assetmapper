from urllib.parse import quote

import requests
from PIL import Image, ImageDraw
from openpyxl import load_workbook
from pptx import Presentation
from pptx import action
from pptx.dml.color import RGBColor
from pptx.enum.action import PP_ACTION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt
from pptx.action import Hyperlink
import os
import requests
import pandas as pd
from urllib.parse import quote
from PyQt5.QtCore import QUrl, Qt, QRect
from PyQt5.QtGui import QColor, QPainter, QPixmap
from PyQt5.QtWidgets import QApplication, QMainWindow, QVBoxLayout, QPushButton, QGraphicsView, QGraphicsScene, QWidget, QGridLayout
from PyQt5.QtWidgets import QGraphicsProxyWidget
from PyQt5.QtWebEngineWidgets import QWebEngineView, QWebEngineSettings


def get_coordinates(api_key, address):
    geocoding_url = 'https://maps.googleapis.com/maps/api/geocode/json'
    params = {
        'address': address,
        'key': api_key
    }
    response = requests.get(geocoding_url, params=params)
    if response.status_code == 200:
        data = response.json()
        if data.get('results'):
            location = data['results'][0].get('geometry', {}).get('location')
            if location:
                latitude = location.get('lat')
                longitude = location.get('lng')
                return latitude, longitude
    return None, None
def create_overlay_widget():
    overlay_widget = QWidget()
    overlay_widget.setStyleSheet("background-color: rgba(255, 0, 0, 0.5); border: 5px solid red;")
    overlay_widget.setFixedSize(int(7 * window.logicalDpiX()), int(4.17 * window.logicalDpiY()))
    return overlay_widget

def get_static_map_image(api_key, latitude, longitude, zoom=16, marker_color='red', marker_size='tiny', nearby_places=[]):
    base_url = 'https://maps.googleapis.com/maps/api/staticmap?'
    params = {
        'center': f'{latitude},{longitude}',
        'zoom': zoom,
        'size': '1280x960',  # Default size for high resolution
        'maptype': 'satellite',
        'markers': f'color:{marker_color}|size:{marker_size}|{latitude},{longitude}',
        'scale': 2,  # Higher scale for higher pixel density
        'key': api_key
    }

    # Add markers for nearby places
    if nearby_places:
        markers = '|'.join([f'label:{place}|{latitude},{longitude}' for place in nearby_places])
        params['markers'] += f'|{markers}'

    response = requests.get(base_url, params=params)
    if response.status_code == 200:
        with open('map_image.png', 'wb') as file:
            file.write(response.content)
        print('Map image saved successfully.')
    else:
        print('Error: Unable to obtain map image.')

def create_border(image, border_thickness):
    draw = ImageDraw.Draw(image)
    width, height = image.size
    draw.rectangle([(0, 0), (width - 1, height - 1)], outline=(0, 0, 0), width=border_thickness)
    return image

def create_hyperlink(prs, text_frame, address):
    p = text_frame.add_paragraph()
    r = p.add_run()
    r.text = address
    hlink = r.hyperlink
    hlink.address = address
    hlink.action = action.Action(pp_action=PP_ACTION.HYPERLINK)
    hlink.action.hyperlink.address = address
    return p


# Specify your API key here
api_key = 'AIzaSyAj7A0lcIbzom4pQJ6_196Fr3-uKHrvc_E'
# Specify the Excel file path
excel_file_path = 'C:/Users/jkitching/OneDrive - Matrix Capital Markets Group/Documents/TestAddys.xlsx'

# Read the Excel file and retrieve the addresses from the specified column
df = pd.read_excel(excel_file_path, header=1)  # Assuming there is no header row
addresses = df.iloc[:, 0].tolist()  # Change the index (0) to the column number containing the addresses

app = QApplication([])

# Create a QMainWindow as the main window
window = QMainWindow()

# Create a QGridLayout for the main layout
layout = QGridLayout()

# Create a QWidget to hold the web view and overlay
widget = QWidget()
widget.setLayout(layout)

# Create a QGraphicsScene to hold the overlay widget
scene = QGraphicsScene()

# Create an overlay widget
overlay_widget = QWidget()
overlay_widget.setStyleSheet("background-color: rgba(255, 0, 0, 0.5); border: 5px solid red;")
overlay_widget.setFixedSize(int(7 * window.logicalDpiX()), int(4.17 * window.logicalDpiY()))

# Center the overlay widget in the middle of the window
overlay_x = (window.width() - overlay_widget.width()) // 2
overlay_y = (window.height() - overlay_widget.height()) // 2
overlay_widget.move(overlay_x, overlay_y)

# Set the overlay widget to be transparent for mouse events
overlay_widget.setAttribute(Qt.WA_TransparentForMouseEvents)

# Add the overlay widget to the scene
scene.addWidget(overlay_widget)

# Create a QGraphicsView to display the scene
graphics_view = QGraphicsView(scene)
graphics_view.setStyleSheet("background: transparent;")
graphics_view.setAttribute(Qt.WA_TransparentForMouseEvents)

# Add the graphics view to the layout
layout.addWidget(graphics_view, 0, 0, 1, 1)


# Create a QPushButton to capture the screenshot
button = QPushButton("Capture Screenshot", window)
button.setFixedHeight(50)  # Set the button height to 50 pixels

# Create a QVBoxLayout for the main layout
main_layout = QVBoxLayout()
main_layout.addWidget(widget)
main_layout.addWidget(button)

# Create a QWidget as the central widget and set the layout
central_widget = QWidget(window)
central_widget.setLayout(main_layout)
window.setCentralWidget(central_widget)

# Initialize the screenshot index
screenshot_index = 0
central_widget.setVisible(True)


def capture_screenshot(webview, overlay, address):
    global screenshot_index

    # Hide the overlay temporarily
    overlay.setVisible(True)

    # Capture the full screenshot
    screenshot = webview.grab()

    # Calculate the crop dimensions
    crop_width = int(7 * webview.logicalDpiX())
    crop_height = int(4.17 * webview.logicalDpiY())
    crop_x = (screenshot.width() - crop_width) // 2
    crop_y = (screenshot.height() - crop_height) // 2

    # Crop the screenshot
    cropped_screenshot = screenshot.copy(QRect(crop_x, crop_y, crop_width, crop_height))

    # Save the cropped screenshot to a file with the address name
    screenshot_path = os.path.join(os.getcwd(), f"{address}.png")
    cropped_screenshot.save(screenshot_path, "PNG")

    screenshot_index += 1

    print("Screenshot saved:", screenshot_path)


def process_next_address():
    if len(addresses) > 0:
        # Get the next address
        address = addresses.pop(0)
        latitude, longitude = get_coordinates(api_key, address)
        # Generate the Google Street View URL for the address
        url = f"https://www.google.com/maps/@?api=1&map_action=pano&viewpoint={latitude},{longitude}&fov=90&heading=235&pitch=10&source=apiv3&hl=en&panoid=ID&sa=X&ved=2ahUKEwiB2NGM68PlAhUDuRoKHQlVDkQQo_oBegQICBAB"

        # Create a QWebEngineView
        view = QWebEngineView()

        # Set web view settings to enable auto resizing
        view.settings().setAttribute(QWebEngineSettings.ShowScrollBars, False)
        view.settings().setAttribute(QWebEngineSettings.JavascriptEnabled, True)
        view.settings().setAttribute(QWebEngineSettings.PluginsEnabled, True)

        # Load the URL
        view.load(QUrl(url))

        # Create a QGraphicsProxyWidget for the overlay widget
        overlay_widget = QGraphicsProxyWidget()
        overlay_widget.setWidget(create_overlay_widget())

        # Add the overlay widget to the scene
        scene.addItem(overlay_widget)

        # Connect the capture_screenshot function to the button's clicked signal
        button.clicked.connect(lambda: capture_screenshot(view, overlay_widget.widget(), address))

        # Add the view to the layout
        layout.addWidget(view, 0, 0, 1, 1)

    else:
        # Close the application when all addresses have been processed
        app.quit()




# Connect the process_next_address function to the button's clicked signal
button.clicked.connect(process_next_address)

window.show()

# Start processing the first address
process_next_address()

app.exec_()

# Specify the distance of the bottom picture from the bottom of the slide (in inches)
bottom_distance_inches = 0.25

# Specify the vertical distance between the two pictures (in inches)
vertical_distance_inches = 0.125

# Specify the thickness of the image border (in inches)
image_border_thickness_inches = 0.03

# Specify the height and width of the rectangle (in inches)
rectangle_height_inches = 0.92
rectangle_width_inches = 7.17

# Specify the distance of the rectangle from the top of the slide (in inches)
rectangle_top_distance_inches = 0.125

# Specify the thickness of the rectangle border (in inches)
rectangle_border_thickness_inches = 0.05

# Specify the dimensions of the second rectangle (in inches)
second_rectangle_height_inches = 0.75
second_rectangle_width_inches = 7

# Specify the position of the second rectangle (left and top coordinates in inches)
second_rectangle_left_inches = 0.25
second_rectangle_top_inches = 0.325

# Specify the border color and width of the second rectangle
second_rectangle_border_color = RGBColor(191, 191, 191)  # Dark gray color
second_rectangle_border_thickness_inches = 0.05


# Specify the font name, size, and style for the address
address_font_name = 'Arial Narrow'
address_font_size = 14
address_font_bold = True

# Specify the font name, size, and style for the data text boxes
data_font_name = 'Arial Narrow'
data_font_size = 14
data_font_bold = True

# Specify the location of the address text (left, top, width, height) in inches
address_box = {
    'left': Inches(0.3),
    'top': Inches(0.23),
    'width': Inches(2.5),
    'height': Inches(1)
}

# Specify the location of the data text boxes (left, top, width, height) in inches
data_b_box = {
    'left': Inches(2.5),
    'top': Inches(.23),
    'width': Inches(2.5),
    'height': Inches(1)
}

data_c_box = {
    'left': Inches(4.7),
    'top': Inches(.23),
    'width': Inches(2.5),
    'height': Inches(1)
}
data_url_box = {
    'left': Inches(8),
    'top': Inches(3),
    'width': Inches(6.9),
    'height': Inches(0.5)
}


# Load the Excel file
workbook = load_workbook(excel_file_path)
sheet = workbook.active

# Create a PowerPoint presentation
presentation = Presentation()

# Set slide dimensions to portrait orientation with a height of 10 inches and width of 7.5 inches
slide_width = Inches(7.5)
slide_height = Inches(10)
presentation.slide_width = slide_width
presentation.slide_height = slide_height

# Iterate through the addresses in the Excel file
for row in sheet.iter_rows(min_row=1, values_only=True):
    address = row[0]
    data_b = row[1]
    data_c = row[2]

    # Specify the desired dimensions for cropping and image size in inches
    desired_width_inches = 7
    desired_height_inches = 4.17

    # Convert inches to pixels based on DPI
    desired_dpi = 96
    desired_width = int(desired_width_inches * desired_dpi)
    desired_height = int(desired_height_inches * desired_dpi)

    # Obtain the coordinates for the address
    latitude, longitude = get_coordinates(api_key, address)
    # Encode the address for the Street View URL
    encoded_address = quote(address, safe='')

    # Generate the Street View URL
    streetview_url = f"https://www.google.com/maps/@?api=1&map_action=pano&viewpoint={latitude},{longitude}&fov=90&heading=235&pitch=10&source=apiv3&hl=en&panoid=ID&sa=X&ved=2ahUKEwiB2NGM68PlAhUDuRoKHQlVDkQQo_oBegQICBAB"
    streetview_url = streetview_url.replace('ID', f'{latitude},{longitude}')
    streetview_url = streetview_url.replace('latitude', str(latitude))
    streetview_url = streetview_url.replace('longitude', str(longitude))
    streetview_url = streetview_url.replace('encoded_address', encoded_address)
    print(f"Street View URL for Slide {streetview_url}")

    if latitude is not None and longitude is not None:
        # Call the function to retrieve the high-resolution map image
        get_static_map_image(api_key, latitude, longitude)

        # Open the original image
        original_image = Image.open('map_image.png')

        # Calculate the cropping dimensions
        left = (original_image.width - desired_width) // 2
        top = (original_image.height - desired_height) // 2
        right = left + desired_width
        bottom = top + desired_height

        # Crop the image
        cropped_image = original_image.crop((left, top, right, bottom))
        cropped_image.save('cropped_map_image.png')

        # Create a slide
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        # Calculate the position and size of the images on the slide
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        # Calculate the total height of the images including the vertical distance and bottom distance
        total_height = (Inches(desired_height_inches) * 2) + Inches(vertical_distance_inches) + Inches(bottom_distance_inches)

        # Calculate the top position of the first image (grey image)
        first_image_top = slide_height - total_height

        # Calculate the left position of the images for horizontal centering
        left = (slide_width - Inches(desired_width_inches)) / 2

        # Convert image border thickness from inches to pixels
        image_border_thickness = int(image_border_thickness_inches * desired_dpi)

        # Create a blank gray image for the first picture with border
        #blank_image = Image.new('RGB', (desired_width, desired_height), color=(192, 192, 192))
        blank_image=Image.open(address+'.png')
        blank_image_with_border = create_border(blank_image, image_border_thickness)
        blank_image_with_border.save('blank_image.png')

        # Add the first image (blank gray image with border) to the slide
        slide.shapes.add_picture('blank_image.png', left, first_image_top, width=Inches(desired_width_inches), height=Inches(desired_height_inches))

        # Create a border for the second image (map image)
        map_image_with_border = create_border(cropped_image, image_border_thickness)
        map_image_with_border.save('map_image_with_border.png')

        # Add the second image (map image with border) to the slide
        second_image_top = first_image_top + Inches(desired_height_inches) + Inches(vertical_distance_inches)
        slide.shapes.add_picture('map_image_with_border.png', left, second_image_top, width=Inches(desired_width_inches), height=Inches(desired_height_inches))

        # Convert rectangle dimensions from inches to pixels
        rectangle_height = int(rectangle_height_inches * desired_dpi)
        rectangle_width = int(rectangle_width_inches * desired_dpi)

        # Calculate the top position of the rectangle
        rectangle_top = first_image_top - Inches(rectangle_top_distance_inches) - Inches(rectangle_height_inches)

        # Convert rectangle border thickness from inches to points
        rectangle_border_thickness = int(rectangle_border_thickness_inches * 72)  # 1 inch = 72 points

        # Add a rectangle shape to the slide
        rectangle_left = (slide_width - Inches(rectangle_width_inches)) / 2
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rectangle_left, rectangle_top, Inches(rectangle_width_inches), Inches(rectangle_height_inches))
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color
        rectangle.line.color.rgb = RGBColor(0, 0, 0)  # Black color
        rectangle.line.width = Pt(rectangle_border_thickness)  # Border thickness in points

        # Convert second rectangle dimensions and positioning from inches to points
        second_rectangle_height = int(second_rectangle_height_inches * 72)  # 1 inch = 72 points
        second_rectangle_width = int(second_rectangle_width_inches * 72)
        second_rectangle_left = Inches(second_rectangle_left_inches)
        second_rectangle_top = Inches(second_rectangle_top_inches)

        # Add the second rectangle shape to the slide
        second_rectangle = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, second_rectangle_left, second_rectangle_top,
            Inches(second_rectangle_width_inches), Inches(second_rectangle_height_inches)
        )
        second_rectangle.fill.background()  # Set the fill to transparent
        second_rectangle.line.color.rgb = second_rectangle_border_color
        second_rectangle.line.width = Pt(second_rectangle_border_thickness_inches)  # Border thickness in points
        second_rectangle.line.dash_style = None  # Remove any dashed line style
        second_rectangle.shadow.inherit = False  # Disable shadows

        # Add the address to the slide
        address_shape = slide.shapes.add_textbox(left=address_box['left'], top=address_box['top'], width=address_box['width'], height=address_box['height'])
        address_text_frame = address_shape.text_frame

        # Set the font properties of the address
        p_address = address_text_frame.add_paragraph()
        p_address.text = address[:address.index(",")]
        p_address.font.name = address_font_name
        p_address.font.size = Pt(address_font_size)
        p_address.font.bold = address_font_bold
        p_address.alignment = PP_ALIGN.LEFT

        # Add the data from column B to the slide
        data_b_shape = slide.shapes.add_textbox(left=data_b_box['left'], top=data_b_box['top'], width=data_b_box['width'], height=data_b_box['height'])
        data_b_text_frame = data_b_shape.text_frame

        # Set the font properties of the data from column B
        p_b = data_b_text_frame.add_paragraph()
        p_b.text = str(data_b)  # Convert data to string
        p_b.font.name = data_font_name
        p_b.font.size = Pt(data_font_size)
        p_b.font.bold = data_font_bold
        p_b.alignment = PP_ALIGN.CENTER

        # Add the data from column C to the slide
        data_c_shape = slide.shapes.add_textbox(left=data_c_box['left'], top=data_c_box['top'], width=data_c_box['width'], height=data_c_box['height'])
        data_c_text_frame = data_c_shape.text_frame

        # Set the font properties of the data from column C
        p_c = data_c_text_frame.add_paragraph()
        p_c.text = str(data_c)  # Convert data to string
        p_c.font.name = data_font_name
        p_c.font.size = Pt(data_font_size)
        p_c.font.bold = data_font_bold
        p_c.alignment = PP_ALIGN.RIGHT

        data_url_shape = slide.shapes.add_textbox(left=data_url_box['left'], top=data_url_box['top'],
                                                  width=data_url_box['width'], height=data_url_box['height'])
        data_url_text_frame = data_url_shape.text_frame

        # Add the Street View URL as a clickable hyperlink
        p_url = data_url_text_frame.add_paragraph()
        r_url = p_url.add_run()
        r_url.text = address
        r_url.font.name = data_font_name
        r_url.font.size = Pt(data_font_size)
        r_url.font.bold = data_font_bold
        hyperlink = r_url.hyperlink
        hyperlink.address = streetview_url

        # Remove the regular text link paragraph
        data_url_text_frame.clear()
        p_url = data_url_text_frame.add_paragraph()
        r_url = p_url.add_run()
        r_url.text = address
        r_url.font.name = data_font_name
        r_url.font.size = Pt(40)
        r_url.font.bold = data_font_bold
        hyperlink = r_url.hyperlink
        hyperlink.address = streetview_url



# Save the PowerPoint presentation
presentation.save('map_slides.pptx')
print('Map slides saved successfully.')
